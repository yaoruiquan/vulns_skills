#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Render a WeChat cover image for vulnerability alert articles."""

from __future__ import annotations

import argparse
import os
import re
import shutil
import subprocess
import sys
import zipfile
from datetime import date
from pathlib import Path


SKILL_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE = SKILL_ROOT / "assets" / "wechat-alert-cover-template.pptx"
CVE_PATTERN = re.compile(r"CVE-\d{4}-\d{4,}", re.IGNORECASE)
LEVEL_PATTERN = re.compile(r"(严重|超危|危急|高危|中危|低危)")


def read_markdown_metadata(markdown_path: Path) -> dict[str, str]:
    text = markdown_path.read_text(encoding="utf-8")
    title = ""
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("# "):
            title = stripped[2:].strip()
            break
    if not title:
        title = markdown_path.stem
    # 去除标题前缀标记如 【已复现】、【风险通告】、【漏洞预警】等
    title = re.sub(r"^【[^】]*】\s*", "", title)
    cve_match = CVE_PATTERN.search(text)
    level_match = LEVEL_PATTERN.search(text)
    summary = ""
    for line in text.splitlines():
        stripped = line.strip()
        if stripped and not stripped.startswith("#") and not stripped.startswith("|"):
            summary = re.sub(r"<[^>]+>", "", stripped)
            break
    return {
        "TITLE": title,
        "CVE": cve_match.group(0).upper() if cve_match else "",
        "LEVEL": level_match.group(1) if level_match else "",
        "DATE": date.today().isoformat(),
        "SUMMARY": summary[:120],
    }


def wrap_text(text: str, limit: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for char in text:
        width = 2 if ord(char) > 127 else 1
        current_width = sum(2 if ord(c) > 127 else 1 for c in current)
        if current and current_width + width > limit:
            lines.append(current)
            current = char
        else:
            current += char
    if current:
        lines.append(current)
    return lines


def wrap_text_by_pixels(text: str, draw, font, max_width: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for char in text:
        candidate = current + char
        bbox = draw.textbbox((0, 0), candidate, font=font)
        if current and bbox[2] - bbox[0] > max_width:
            lines.append(current)
            current = char
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines


def option_value(value: str) -> bool:
    return str(value or "").strip().lower() in {"1", "true", "yes", "y", "on", "是", "有", "已公开", "已复现"}


def first_pptx_background(template: Path) -> Path | None:
    """Extract the first media image from a PPTX template as a clean background."""
    try:
        with zipfile.ZipFile(template) as archive:
            media = [
                name
                for name in archive.namelist()
                if name.startswith("ppt/media/") and name.lower().endswith((".png", ".jpg", ".jpeg"))
            ]
            if not media:
                return None
            target = Path(os.environ.get("WECHAT_COVER_OUTPUT_DIR", "/tmp/md2wechat-covers")) / f"{template.stem}-background{Path(media[0]).suffix}"
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_bytes(archive.read(media[0]))
            return target
    except zipfile.BadZipFile:
        return None


def draw_template_cover(
    template_image: Path,
    output: Path,
    values: dict[str, str],
    clear_existing: bool = False,
    font_path: str = "",
) -> Path:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as exc:
        raise SystemExit("缺少 Pillow，先执行：python3 -m pip install Pillow") from exc

    image = Image.open(template_image).convert("RGB")
    # WeChat cover images are commonly 900x383. Keep the provided template ratio.
    if image.size != (900, 383):
        image = image.resize((900, 383))
    draw = ImageDraw.Draw(image)
    width, height = image.size

    def font(size: int):
        if font_path and Path(font_path).is_file():
            return ImageFont.truetype(font_path, size=size)
        candidates = [
            "/Users/yao/Library/Fonts/SourceHanSansCN-Bold.otf",
            "/Users/yao/Library/Fonts/SourceHanSansSC-Bold.otf",
            "/Library/Fonts/SourceHanSansCN-Bold.otf",
            "/Library/Fonts/SourceHanSansSC-Bold.otf",
            "/System/Library/Fonts/STHeiti Medium.ttc",
            "/System/Library/Fonts/Hiragino Sans GB.ttc",
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/System/Library/Fonts/PingFang.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
        ]
        for path in candidates:
            if Path(path).is_file():
                try:
                    return ImageFont.truetype(path, size=size)
                except OSError:
                    continue
        return ImageFont.load_default()

    option_font = font(17)

    if clear_existing:
        draw.rounded_rectangle([54, 92, 842, 265], radius=8, fill=(244, 250, 250))

    title_box_width = 770
    title_font_size = 36
    title_font = font(title_font_size)
    title_lines = wrap_text_by_pixels(values["TITLE"], draw, title_font, title_box_width)[:3]
    line_height = 48
    y = 150
    for line in title_lines:
        draw.text((70, y), line, fill="#34294B", font=title_font)
        y += line_height

    line_y = max(220, y + 2)
    draw.line([70, line_y, 760, line_y], fill="#34294B", width=2)

    options = [
        ("POC", values.get("POC", "")),
        ("EXP", values.get("EXP", "")),
        ("在野利用", values.get("WILD", "")),
        ("研究情况", values.get("RESEARCH", "")),
    ]
    x = 70
    for label, enabled in options:
        checked = option_value(enabled)
        option_y = line_y + 23
        box = [x, option_y, x + 15, option_y + 15]
        draw.rectangle(box, outline="#34294B", width=1, fill="#BFEFD6" if checked else None)
        bbox = draw.textbbox((0, 0), label, font=option_font)
        text_h = bbox[3] - bbox[1]
        text_y = option_y + (15 - text_h) // 2
        draw.text((x + 28, text_y), label, fill="#34294B", font=option_font)
        x += 83 if label in {"POC", "EXP"} else 112

    output.parent.mkdir(parents=True, exist_ok=True)
    image.save(output)
    return output


def render_png_template(template: Path, output: Path, values: dict[str, str], font_path: str = "") -> Path:
    return draw_template_cover(template, output, values, clear_existing=True, font_path=font_path)


def fill_pptx_template(template: Path, output_pptx: Path, values: dict[str, str]) -> Path:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("缺少 python-pptx，先执行：python3 -m pip install python-pptx") from exc

    presentation = Presentation(str(template))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    for key, value in values.items():
                        text = text.replace("{{" + key + "}}", value)
                    run.text = text
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_pptx))
    return output_pptx


def export_pptx_to_png(pptx_path: Path, output_png: Path) -> Path | None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice") or "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if not Path(soffice).exists() and not shutil.which(str(soffice)):
        return None
    output_dir = output_png.parent
    subprocess.run(
        [str(soffice), "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_path)],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    pdf_path = output_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.is_file():
        return None
    if shutil.which("pdftoppm"):
        subprocess.run(["pdftoppm", "-png", "-singlefile", str(pdf_path), str(output_png.with_suffix(""))], check=True)
        return output_png if output_png.is_file() else None
    if shutil.which("sips"):
        subprocess.run(["sips", "-s", "format", "png", str(pdf_path), "--out", str(output_png)], check=True, stdout=subprocess.DEVNULL)
        return output_png if output_png.is_file() else None
    return None


def main() -> int:
    parser = argparse.ArgumentParser(description="Render WeChat cover from alert Markdown and PNG/PPTX template")
    parser.add_argument("markdown")
    parser.add_argument("--template", default=os.environ.get("WECHAT_COVER_TEMPLATE", str(DEFAULT_TEMPLATE)))
    parser.add_argument("--output", default="")
    parser.add_argument("--title", default="")
    parser.add_argument("--cve", default="")
    parser.add_argument("--level", default="")
    parser.add_argument("--summary", default="")
    parser.add_argument("--poc", default="")
    parser.add_argument("--exp", default="")
    parser.add_argument("--wild", default="")
    parser.add_argument("--research", default="")
    parser.add_argument("--font", default=os.environ.get("WECHAT_COVER_FONT", ""), help="思源黑体等本地字体文件路径")
    args = parser.parse_args()

    markdown_path = Path(args.markdown).expanduser().resolve()
    template = Path(args.template).expanduser().resolve() if args.template else None
    if not markdown_path.is_file():
        raise SystemExit(f"Markdown 不存在: {markdown_path}")
    if template is None or not template.is_file():
        raise SystemExit("封面模板不存在，请通过 --template 或 WECHAT_COVER_TEMPLATE 指定 .png/.pptx")

    values = read_markdown_metadata(markdown_path)
    for key, value in {
        "TITLE": args.title,
        "CVE": args.cve,
        "LEVEL": args.level,
        "SUMMARY": args.summary,
        "POC": args.poc,
        "EXP": args.exp,
        "WILD": args.wild,
        "RESEARCH": args.research,
    }.items():
        if value:
            values[key] = value

    output = Path(args.output).expanduser() if args.output else Path(os.environ.get("WECHAT_COVER_OUTPUT_DIR", "/tmp/md2wechat-covers")) / f"{markdown_path.stem}-cover.png"
    output = output.resolve()

    if template.suffix.lower() == ".png":
        result = render_png_template(template, output, values, font_path=args.font)
        print(result)
        return 0
    if template.suffix.lower() == ".pptx":
        background = first_pptx_background(template)
        if background:
            result = draw_template_cover(background, output, values, clear_existing=False, font_path=args.font)
            print(result)
            return 0
        filled_pptx = output.with_suffix(".pptx")
        fill_pptx_template(template, filled_pptx, values)
        exported = export_pptx_to_png(filled_pptx, output)
        if exported:
            print(exported)
            return 0
        print(f"已生成 PPTX，但未能自动导出 PNG: {filled_pptx}", file=sys.stderr)
        print("请安装 LibreOffice 和 pdftoppm，或手动导出 PNG 后作为 --cover 使用。", file=sys.stderr)
        return 2
    raise SystemExit("封面模板仅支持 .png 或 .pptx")


if __name__ == "__main__":
    raise SystemExit(main())
